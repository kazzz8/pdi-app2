#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
現場作業管理システム（仮称）システム概要説明資料
PowerPoint生成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUT  = os.path.join(os.path.dirname(os.path.abspath(__file__)), "system_overview.pptx")
NAVY       = RGBColor(0x1A, 0x3A, 0x6B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x1E, 0x8B, 0x4C)
TEAL       = RGBColor(0x0E, 0x6E, 0x75)
ORANGE     = RGBColor(0xD0, 0x6A, 0x00)
LIGHT_BLUE = RGBColor(0xEA, 0xF0, 0xFB)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF8)
MID_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
FONT = "游ゴシック"
W = Inches(13.33)
H = Inches(7.5)


# ============================================================
# ヘルパー
# ============================================================
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, lc=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = lc
        if lw: s.line.width = Pt(lw)
    else:
        s.line.color.rgb = fill or WHITE
    return s

def para(tf, text, fs=11, bold=False, color=DARK_GRAY,
         align=PP_ALIGN.LEFT, first=True):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(fs)
    r.font.bold = bold; r.font.name = FONT
    r.font.color.rgb = color
    return p

def mlbox(slide, text, l, t, w, h,
          fs=11, bold=False, color=DARK_GRAY,
          align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = valign
    for i, ln in enumerate(text.split("\n")):
        para(tf, ln, fs=fs, bold=bold, color=color, align=align, first=(i == 0))
    return box

def lrect(slide, text, l, t, w, h, fill, tc=WHITE,
          fs=12, bold=True, align=PP_ALIGN.CENTER,
          valign=MSO_ANCHOR.MIDDLE, lc=None, lw=None):
    s = rect(slide, l, t, w, h, fill=fill, lc=lc or fill, lw=lw)
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = valign
    para(tf, text, fs=fs, bold=bold, color=tc, align=align)
    return s

def title_bar(slide, text):
    lrect(slide, text,
          Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.65),
          fill=NAVY, tc=WHITE, fs=17, bold=True, align=PP_ALIGN.LEFT)


# ============================================================
# スライド1: タイトル
# ============================================================
def s1(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=NAVY)
    rect(sl, Inches(0.55), Inches(2.4), Inches(0.1), Inches(2.5), fill=WHITE)
    mlbox(sl, "現場作業管理システム（仮称）\nシステム概要説明資料",
          Inches(0.85), Inches(2.2), Inches(11.5), Inches(2.0),
          fs=26, bold=True, color=WHITE)
    mlbox(sl, "機能・画面構成・技術スタック・インフラ",
          Inches(0.85), Inches(4.35), Inches(11.0), Inches(0.6),
          fs=13, color=RGBColor(0xBB, 0xCC, 0xEE))
    mlbox(sl, "2026年4月　情報システム部",
          Inches(0.85), Inches(5.2), Inches(6), Inches(0.5),
          fs=11, color=RGBColor(0x99, 0xAA, 0xCC))


# ============================================================
# スライド2: システム概要
# ============================================================
def s2(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "１．システム概要")

    items = [
        (NAVY,   "目的",     "完成車のPDI（出荷前整備）作業における記録・報告業務を\nスマートフォン・タブレット1台で完結させる"),
        (GREEN,  "対象業務", "車両入庫から出庫までの全作業工程（14種類）の\n作業時間記録・点検報告・不具合管理"),
        (TEAL,   "利用者",   "現場作業者・班長・管理者　計 200〜300名"),
        (ORANGE, "動作環境", "社内ネットワーク内のWindows Serverで稼働する\nPWA（Progressive Web App）。専用アプリのインストール不要。"),
    ]

    iw, ih = Inches(5.85), Inches(1.4)
    positions = [
        (Inches(0.5),  Inches(1.25)),
        (Inches(6.9),  Inches(1.25)),
        (Inches(0.5),  Inches(2.85)),
        (Inches(6.9),  Inches(2.85)),
    ]

    for i, (color, label, body) in enumerate(items):
        l, t = positions[i]
        rect(sl, l, t, iw, ih, fill=LIGHT_GRAY, lc=color, lw=1.0)
        lrect(sl, label, l, t, Inches(1.5), ih, fill=color,
              tc=WHITE, fs=12, bold=True)
        mlbox(sl, body, l + Inches(1.6), t + Inches(0.2),
              iw - Inches(1.7), ih - Inches(0.3),
              fs=11, color=DARK_GRAY, valign=MSO_ANCHOR.MIDDLE)

    # PWAの補足
    note = rect(sl, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.7),
                fill=RGBColor(0xFF, 0xFB, 0xE6),
                lc=RGBColor(0xE8, 0xB0, 0x00), lw=0.75)
    para(note.text_frame,
         "PWAとは：Webブラウザで動作しながらスマートフォンのネイティブアプリに近い操作感を実現する技術。"
         "App Storeからのインストール不要で、ブラウザのアドレスバーからアクセスするだけで利用できる。",
         fs=10, color=RGBColor(0x7D, 0x5A, 0x00))
    note.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    note.text_frame.word_wrap = True


# ============================================================
# スライド3: 対応業務工程
# ============================================================
def s3(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "２．対応する業務工程（14工程）")

    processes = [
        ("積み込み・降ろし", "LOADING",        NAVY),
        ("保管",            "STORAGE",        NAVY),
        ("回送",            "RETRIEVAL",      NAVY),
        ("洗車",            "WASH",           TEAL),
        ("磨き",            "POLISH",         TEAL),
        ("L点検",           "L_INSPECTION",   GREEN),
        ("S点検",           "S_INSPECTION",   GREEN),
        ("補修",            "REPAIR",         RED),
        ("防錆",            "RUST_PREVENTION",TEAL),
        ("ドレスアップ部品払い出し", "PARTS_ISSUE", ORANGE),
        ("ドレスアップ",    "DRESS_UP",       ORANGE),
        ("コーティング",    "COATING",        TEAL),
        ("完成検査",        "FINAL_INSPECTION",GREEN),
        ("その他（計画外）","OTHER",          MID_GRAY),
    ]

    cols = 4
    bw, bh = Inches(2.9), Inches(0.85)
    ml, mt = Inches(0.5), Inches(1.25)
    gap_x, gap_y = Inches(0.3), Inches(0.22)

    for i, (label, code, color) in enumerate(processes):
        col = i % cols
        row = i // cols
        l = ml + col * (bw + gap_x)
        t = mt + row * (bh + gap_y)

        rect(sl, l, t, bw, bh, fill=LIGHT_GRAY, lc=color, lw=0.75)
        rect(sl, l, t, Inches(0.28), bh, fill=color)

        mlbox(sl, label, l + Inches(0.38), t + Inches(0.08),
              bw - Inches(0.5), Inches(0.45),
              fs=11, bold=True, color=DARK_GRAY, valign=MSO_ANCHOR.MIDDLE)
        mlbox(sl, code, l + Inches(0.38), t + Inches(0.52),
              bw - Inches(0.5), Inches(0.28),
              fs=8.5, color=RGBColor(0x88, 0x88, 0x88))

    # 凡例
    legend = [
        (NAVY,   "入庫・保管・搬送"),
        (TEAL,   "整備・防錆・コーティング"),
        (GREEN,  "点検・完成検査"),
        (RED,    "補修"),
        (ORANGE, "ドレスアップ"),
    ]
    lx = Inches(0.5)
    lt = Inches(6.55)
    for j, (color, label) in enumerate(legend):
        rect(sl, lx + j * Inches(2.45), lt, Inches(0.18), Inches(0.22), fill=color)
        mlbox(sl, label, lx + j * Inches(2.45) + Inches(0.25), lt,
              Inches(2.1), Inches(0.25), fs=9, color=RGBColor(0x55, 0x55, 0x55))


# ============================================================
# スライド4: 利用者とロール
# ============================================================
def s4(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "３．利用者とロール（権限区分）")

    roles = [
        (GREEN,  "一般作業者\nGENERAL",
         [
             "本日の作業一覧を確認する",
             "バーコードスキャンで車両を確認して作業を開始する",
             "作業タイマーを操作する（開始・中断・再開・完了）",
             "L点検・S点検の不具合を展開図に記録する",
             "不具合箇所の写真を撮影・登録する",
             "計画外作業を自ら開始する",
         ],
         "最も多い利用者層。現場での操作をスマホで完結。"),
        (ORANGE, "班長\nTEAM_LEADER",
         [
             "一般作業者の全機能",
             "チームメンバーの作業進捗をリアルタイムで確認する",
             "作業計画の調整・割り当て変更を行う",
             "不具合の状況を集約して確認する",
         ],
         "チーム単位の進捗管理・作業指示を担当。"),
        (NAVY,   "管理者\nMANAGER",
         [
             "班長の全機能",
             "ユーザーアカウントを作成・管理する",
             "作業計画データをCSV等でインポートする",
             "全チームの作業実績・KPIを集計・出力する",
             "システム設定・マスタデータを管理する",
         ],
         "全体管理・データ分析・システム運用を担当。"),
    ]

    rw, rh = Inches(3.85), Inches(5.2)
    lefts = [Inches(0.45), Inches(4.6), Inches(8.75)]
    top = Inches(1.2)

    for i, (color, title, perms, desc) in enumerate(roles):
        l = lefts[i]
        rect(sl, l, top, rw, rh, fill=LIGHT_GRAY, lc=color, lw=1.0)
        lrect(sl, title, l, top, rw, Inches(0.9),
              fill=color, tc=WHITE, fs=13, bold=True)

        for j, perm in enumerate(perms):
            pt = top + Inches(1.05) + j * Inches(0.58)
            rect(sl, l + Inches(0.18), pt + Inches(0.15),
                 Inches(0.12), Inches(0.12), fill=color)
            mlbox(sl, perm,
                  l + Inches(0.4), pt,
                  rw - Inches(0.55), Inches(0.55),
                  fs=9.5, color=DARK_GRAY)

        # 役割説明
        desc_box = rect(sl, l + Inches(0.15),
                        top + rh - Inches(0.8),
                        rw - Inches(0.3), Inches(0.65),
                        fill=WHITE, lc=color, lw=0.5)
        para(desc_box.text_frame, desc, fs=9, color=color, bold=True)
        desc_box.text_frame.word_wrap = True
        desc_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


# ============================================================
# スライド5: 画面フロー
# ============================================================
def s5(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "４．画面フロー")

    screens = [
        (NAVY,  "ログイン画面",       "社員番号＋\nパスワード入力"),
        (GREEN, "ダッシュボード",      "本日の作業一覧\nステータス確認"),
        (TEAL,  "作業開始確認",        "バーコードスキャン\nで車両照合"),
        (ORANGE,"作業中画面",          "タイマー計測\n展開図・不具合記録"),
        (GREEN, "完了画面",            "点検報告書\n確認・保存"),
    ]

    bw, bh = Inches(2.1), Inches(2.0)
    lt = Inches(2.2)
    gap = Inches(0.38)

    for i, (color, title, desc) in enumerate(screens):
        l = Inches(0.45) + i * (bw + gap)

        # 画面ボックス（スマホ風）
        rect(sl, l, lt, bw, bh, fill=LIGHT_GRAY, lc=color, lw=1.2)
        lrect(sl, title, l, lt, bw, Inches(0.55), fill=color,
              tc=WHITE, fs=11, bold=True)
        mlbox(sl, desc, l + Inches(0.1), lt + Inches(0.65),
              bw - Inches(0.2), Inches(1.2),
              fs=10, color=DARK_GRAY, align=PP_ALIGN.CENTER,
              valign=MSO_ANCHOR.MIDDLE)

        # 矢印
        if i < 4:
            ax = l + bw + Inches(0.05)
            mlbox(sl, "→", ax, lt + bh / 2 - Inches(0.25),
                  Inches(0.28), Inches(0.5),
                  fs=18, bold=True, color=NAVY,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 中断フロー
    pause_top = Inches(4.45)
    rect(sl, Inches(6.0), pause_top, Inches(3.0), Inches(1.7),
         fill=RGBColor(0xFF, 0xFB, 0xE6), lc=ORANGE, lw=0.75)
    lrect(sl, "一時中断",
          Inches(6.0), pause_top, Inches(3.0), Inches(0.5),
          fill=ORANGE, tc=WHITE, fs=11, bold=True)
    mlbox(sl, "中断理由を選択・入力\nダッシュボードに戻る\n↓\n再開ボタンから作業を継続",
          Inches(6.1), pause_top + Inches(0.55), Inches(2.8), Inches(1.0),
          fs=9.5, color=DARK_GRAY)

    # 中断矢印
    mlbox(sl, "↕ 中断/再開",
          Inches(5.5), Inches(3.8), Inches(1.2), Inches(0.35),
          fs=9, color=ORANGE, align=PP_ALIGN.CENTER)

    # ルート注記
    mlbox(sl, "※ 計画外作業はダッシュボードから直接「作業中画面」に遷移",
          Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.35),
          fs=9, color=RGBColor(0x88, 0x88, 0x88))


# ============================================================
# スライド6: 主要機能 ① ダッシュボード・作業管理
# ============================================================
def s6(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "５．主要機能 ①　ダッシュボード・作業管理")

    # 左: ダッシュボード
    lrect(sl, "ダッシュボード",
          Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.5),
          fill=GREEN, tc=WHITE, fs=12, bold=True)

    dash_items = [
        "本日割り当てられた作業計画を一覧表示",
        "各作業のステータスをリアルタイムで表示",
        "計画外作業をその場で開始できる",
        "中断中の作業がある場合はバナーで通知",
    ]
    status_map = [
        (GREEN,                        "完了",     "作業が完了"),
        (RGBColor(0x25, 0x6B, 0xBE),   "作業中",   "現在進行中"),
        (RGBColor(0xCA, 0x85, 0x00),   "中断中",   "一時停止中"),
        (RED,                           "遅延",     "計画時間超過"),
        (RGBColor(0x55, 0x55, 0x55),   "予定通り", "未着手・計画内"),
    ]

    for j, item in enumerate(dash_items):
        t = Inches(1.85) + j * Inches(0.52)
        rect(sl, Inches(0.65), t + Inches(0.1), Inches(0.1), Inches(0.1), fill=GREEN)
        mlbox(sl, item, Inches(0.85), t, Inches(5.3), Inches(0.45),
              fs=10.5, color=DARK_GRAY)

    mlbox(sl, "ステータス表示",
          Inches(0.65), Inches(4.05), Inches(5.0), Inches(0.35),
          fs=10, bold=True, color=DARK_GRAY)
    for j, (color, label, desc) in enumerate(status_map):
        st = Inches(4.5) + j * Inches(0.44)
        rect(sl, Inches(0.65), st, Inches(0.85), Inches(0.3), fill=color, lc=color)
        para_box = sl.shapes.add_textbox(Inches(0.65), st, Inches(0.85), Inches(0.3))
        tf = para_box.text_frame
        para(tf, label, fs=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        mlbox(sl, desc, Inches(1.6), st, Inches(4.3), Inches(0.3),
              fs=9.5, color=DARK_GRAY, valign=MSO_ANCHOR.MIDDLE)

    # 区切り線
    rect(sl, Inches(6.6), Inches(1.15), Inches(0.03), Inches(5.5), fill=MID_GRAY)

    # 右: バーコードスキャン
    lrect(sl, "車両バーコードスキャン（作業開始時）",
          Inches(6.9), Inches(1.2), Inches(5.9), Inches(0.5),
          fill=TEAL, tc=WHITE, fs=12, bold=True)

    scan_items = [
        ("目的",     "作業対象の車両と作業計画が一致することを確認する"),
        ("操作",     "作業開始ボタンを押す → カメラが起動 → バーコードにかざす"),
        ("判定",     "スキャン結果と計画の整理番号が一致した場合のみ\n「作業開始」ボタンが活性化する"),
        ("不一致時", "エラー表示。再スキャンを促す。誤った車両への作業を防止する"),
    ]
    for j, (label, body) in enumerate(scan_items):
        t = Inches(1.85) + j * Inches(0.88)
        lrect(sl, label, Inches(6.9), t, Inches(1.1), Inches(0.55),
              fill=TEAL, tc=WHITE, fs=10, bold=True)
        mlbox(sl, body, Inches(8.1), t + Inches(0.05),
              Inches(4.6), Inches(0.7), fs=10, color=DARK_GRAY)


# ============================================================
# スライド7: 主要機能 ② 作業タイマー・進捗管理
# ============================================================
def s7(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "５．主要機能 ②　作業タイマー・進捗管理")

    # 左: タイマー
    lrect(sl, "リアルタイム作業タイマー",
          Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.5),
          fill=NAVY, tc=WHITE, fs=12, bold=True)

    timer_items = [
        ("計測方式", "作業開始と同時にタイマーが自動スタート。\n一時中断・再開の累積時間を自動計算する。"),
        ("表示",     "作業中画面に「MM:SS」形式で経過時間を大きく表示。"),
        ("超過警告", "標準作業時間を超過した場合、タイマーが赤色に変化し\n「超過中」の警告を表示する。"),
        ("データ保存","完了時に開始時刻・終了時刻・累積作業時間をDBに記録。"),
    ]
    for j, (label, body) in enumerate(timer_items):
        t = Inches(1.85) + j * Inches(1.1)
        lrect(sl, label, Inches(0.5), t, Inches(1.3), Inches(0.65),
              fill=NAVY, tc=WHITE, fs=10, bold=True)
        mlbox(sl, body, Inches(1.9), t + Inches(0.05),
              Inches(4.3), Inches(0.8), fs=10, color=DARK_GRAY)

    # 区切り線
    rect(sl, Inches(6.6), Inches(1.15), Inches(0.03), Inches(5.5), fill=MID_GRAY)

    # 右: 中断・再開
    lrect(sl, "一時中断・再開機能",
          Inches(6.9), Inches(1.2), Inches(5.9), Inches(0.5),
          fill=ORANGE, tc=WHITE, fs=12, bold=True)

    pause_items = [
        ("中断理由", "班長対応・緊急補修・部品待ち・休憩・その他\n（定型選択 or 自由入力）"),
        ("再開",     "ダッシュボードの「再開」ボタンから前回の続きで\n作業を再開できる。累積時間は引き継がれる。"),
        ("同時中断", "複数の中断中作業を持てないよう制御されている。\n（誤操作・二重記録の防止）"),
        ("記録",     "中断理由・中断時刻・再開時刻をすべてDBに記録。\n班長・管理者が後から確認できる。"),
    ]
    for j, (label, body) in enumerate(pause_items):
        t = Inches(1.85) + j * Inches(1.1)
        lrect(sl, label, Inches(6.9), t, Inches(1.3), Inches(0.65),
              fill=ORANGE, tc=WHITE, fs=10, bold=True)
        mlbox(sl, body, Inches(8.3), t + Inches(0.05),
              Inches(4.4), Inches(0.8), fs=10, color=DARK_GRAY)


# ============================================================
# スライド8: 主要機能 ③ 不具合記録
# ============================================================
def s8(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "５．主要機能 ③　不具合記録（L点検・S点検）")

    # 左: 展開図
    lrect(sl, "車両展開図タップ入力",
          Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.5),
          fill=RED, tc=WHITE, fs=12, bold=True)

    diag_items = [
        "車両の展開図（上面・側面・前後）を画面に表示",
        "不具合箇所をタップしてピンを立てる",
        "ピンはA/B/C程度ごとに色分けして表示（黄・橙・赤）",
        "ピンをタップすると記録済みの不具合詳細が表示される",
        "不具合は何件でも登録可能",
    ]
    for j, item in enumerate(diag_items):
        t = Inches(1.85) + j * Inches(0.55)
        rect(sl, Inches(0.65), t + Inches(0.12), Inches(0.1), Inches(0.1), fill=RED)
        mlbox(sl, item, Inches(0.85), t, Inches(5.3), Inches(0.5),
              fs=10.5, color=DARK_GRAY)

    # 展開図イメージ（テキストで代用）
    diag_mock = rect(sl, Inches(0.65), Inches(4.6), Inches(5.5), Inches(1.6),
                     fill=RGBColor(0xF8, 0xF8, 0xF8), lc=MID_GRAY, lw=0.5)
    mlbox(sl, "【展開図イメージ】\n前面 ／ 左側面 ／ 後面 ／ 右側面 ／ 上面 を1画面で表示\nタップした座標をX%・Y%で記録",
          Inches(0.85), Inches(4.7), Inches(5.1), Inches(1.3),
          fs=9.5, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER,
          valign=MSO_ANCHOR.MIDDLE)

    # 区切り線
    rect(sl, Inches(6.6), Inches(1.15), Inches(0.03), Inches(5.5), fill=MID_GRAY)

    # 右: 不具合記録項目
    lrect(sl, "不具合記録項目",
          Inches(6.9), Inches(1.2), Inches(5.9), Inches(0.5),
          fill=RED, tc=WHITE, fs=12, bold=True)

    fields = [
        ("不具合種類",    "線傷 / 面傷 / 塗装 / へこみ / 欠け / 汚れ / その他\n（ボタン選択）"),
        ("程度",         "A（軽微）/ B（中程度）/ C（重大）\n（ボタン選択。展開図ピンの色に反映される）"),
        ("補修見込み時間","10・30・60・90・120分のプリセット + 5分刻みの手動調整"),
        ("写真",         "スマホカメラで撮影（任意）。JPEG圧縮してサーバーに保存。"),
        ("座標",         "展開図上のタップ位置をX%・Y%で自動記録"),
    ]
    for j, (label, body) in enumerate(fields):
        t = Inches(1.85) + j * Inches(0.98)
        lrect(sl, label, Inches(6.9), t, Inches(2.0), Inches(0.6),
              fill=RGBColor(0xFF, 0xEE, 0xEE), tc=RED, fs=10, bold=True,
              lc=RED, lw=0.5)
        mlbox(sl, body, Inches(9.0), t + Inches(0.05),
              Inches(3.75), Inches(0.75), fs=10, color=DARK_GRAY)


# ============================================================
# スライド9: 技術スタック
# ============================================================
def s9(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "６．技術スタック")

    stacks = [
        (NAVY,   "フロントエンド",
         [("フレームワーク", "Next.js 15（App Router）"),
          ("言語",          "TypeScript"),
          ("スタイル",       "Tailwind CSS"),
          ("PWA対応",        "manifest.json・Service Worker"),
          ("バーコード",      "カメラAPI（ブラウザ標準）"),]),
        (TEAL,   "バックエンド",
         [("API",            "Next.js API Routes（サーバーサイド）"),
          ("認証",            "NextAuth.js + bcrypt"),
          ("ORM",             "Prisma（TypeScript対応）"),
          ("セッション",       "JWT（サーバーサイドで管理）"),
          ("ファイル",         "写真アップロードAPI（サーバー保存）"),]),
        (GREEN,  "データベース",
         [("DBMS",           "PostgreSQL（移行時にSQL Serverも対応可）"),
          ("主要テーブル",    "Vehicle / WorkPlan / WorkLog\nInspectionReport / Defect / User"),
          ("マイグレーション","Prisma Migrate（バージョン管理済み）"),
          ("シード",          "初期データ投入スクリプト完備"),]),
        (ORANGE, "インフラ（予定）",
         [("OS",             "Windows Server"),
          ("Webサーバー",     "IIS（リバースプロキシ）"),
          ("プロセス管理",    "PM2 または NSSM"),
          ("Node.js",        "v20 LTS"),
          ("SSL",            "社内CA発行証明書"),]),
    ]

    sw, sh = Inches(2.8), Inches(4.6)
    lefts = [Inches(0.45), Inches(3.55), Inches(6.65), Inches(9.75)]
    top = Inches(1.2)

    for i, (color, title, items) in enumerate(stacks):
        l = lefts[i]
        rect(sl, l, top, sw, sh, fill=LIGHT_GRAY, lc=color, lw=1.0)
        lrect(sl, title, l, top, sw, Inches(0.5), fill=color, tc=WHITE, fs=12)

        for j, (key, val) in enumerate(items):
            t = top + Inches(0.6) + j * Inches(0.78)
            mlbox(sl, key, l + Inches(0.15), t,
                  sw - Inches(0.25), Inches(0.3),
                  fs=9, bold=True, color=color)
            mlbox(sl, val, l + Inches(0.15), t + Inches(0.3),
                  sw - Inches(0.25), Inches(0.42),
                  fs=9.5, color=DARK_GRAY)
            if j < len(items) - 1:
                rect(sl, l + Inches(0.1), t + Inches(0.72),
                     sw - Inches(0.2), Inches(0.01), fill=MID_GRAY)


# ============================================================
# スライド10: システム構成
# ============================================================
def s10(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "７．システム構成（Windows Server 構成案）")

    # Windows Server ボックス全体
    rect(sl, Inches(3.5), Inches(1.4), Inches(8.8), Inches(4.9),
         fill=RGBColor(0xF0, 0xF4, 0xF8), lc=NAVY, lw=1.5)
    mlbox(sl, "Windows Server（社内）",
          Inches(3.7), Inches(1.45), Inches(4.0), Inches(0.4),
          fs=10, bold=True, color=NAVY)

    # クライアント
    lrect(sl, "クライアント\n（スマホ/タブレット）",
          Inches(0.35), Inches(3.0), Inches(2.7), Inches(1.4),
          fill=GREEN, tc=WHITE, fs=11, bold=True)

    # 矢印: クライアント → IIS
    mlbox(sl, "HTTPS\n社内LAN",
          Inches(3.1), Inches(3.3), Inches(0.6), Inches(0.6),
          fs=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    mlbox(sl, "→",
          Inches(3.1), Inches(3.5), Inches(0.5), Inches(0.5),
          fs=20, bold=True, color=NAVY,
          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # IIS
    lrect(sl, "IIS\n（リバースプロキシ・SSL終端）",
          Inches(3.7), Inches(2.5), Inches(3.5), Inches(1.2),
          fill=NAVY, tc=WHITE, fs=10, bold=True)

    # 矢印: IIS → Node.js
    mlbox(sl, "↓",
          Inches(5.3), Inches(3.8), Inches(0.5), Inches(0.45),
          fs=18, bold=True, color=NAVY,
          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # Node.js
    lrect(sl, "Node.js + Next.js\n（アプリケーションサーバー）",
          Inches(3.7), Inches(4.3), Inches(3.5), Inches(1.2),
          fill=TEAL, tc=WHITE, fs=10, bold=True)
    mlbox(sl, "PM2 または NSSM でサービス化",
          Inches(3.75), Inches(5.5), Inches(3.4), Inches(0.35),
          fs=9, color=RGBColor(0x55, 0x55, 0x55))

    # 矢印: Node → DB
    mlbox(sl, "→",
          Inches(7.4), Inches(4.75), Inches(0.45), Inches(0.45),
          fs=18, bold=True, color=NAVY,
          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # DB
    lrect(sl, "PostgreSQL\n（データベース）",
          Inches(7.9), Inches(4.3), Inches(2.8), Inches(1.2),
          fill=GREEN, tc=WHITE, fs=10, bold=True)

    # ファイルストレージ
    lrect(sl, "ファイルストレージ\n（写真保存）",
          Inches(7.9), Inches(2.5), Inches(2.8), Inches(1.2),
          fill=ORANGE, tc=WHITE, fs=10, bold=True)
    mlbox(sl, "\\server\\share\\photos\\ など",
          Inches(8.0), Inches(3.75), Inches(2.5), Inches(0.35),
          fs=9, color=RGBColor(0x55, 0x55, 0x55))

    # 矢印: Node → ファイル
    mlbox(sl, "→",
          Inches(7.4), Inches(2.95), Inches(0.45), Inches(0.45),
          fs=18, bold=True, color=NAVY,
          align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 注記
    mlbox(sl, "※ DBはSQL Serverへの変更も対応可（Prismaのprovider設定変更のみ）\n"
              "※ 社内Active Directoryと連携したSSO認証への移行も将来的に対応可能",
          Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6),
          fs=9, color=RGBColor(0x66, 0x66, 0x66))


# ============================================================
# スライド11: データ管理
# ============================================================
def s11(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "８．データ管理（主要テーブルと情報の流れ）")

    tables = [
        (NAVY,   "User",             "社員番号・氏名・ロール・\n班情報を管理"),
        (NAVY,   "Team",             "班の名称と\nメンバー構成"),
        (TEAL,   "Vehicle",          "整理番号・バーコード・\n車種・色・点検区分・状態"),
        (TEAL,   "WorkPlan",         "誰がいつどの車両に\n何の作業をするかの計画"),
        (GREEN,  "WorkLog",          "作業の開始・中断・\n再開・完了の実績記録"),
        (GREEN,  "InspectionReport", "点検報告書\n（点検種別・時間）"),
        (RED,    "Defect",           "不具合の箇所・種類・\n程度・写真・補修状況"),
        (ORANGE, "RepairLog",        "補修作業の実績"),
        (ORANGE, "Photo",            "写真ファイルパスと\n撮影者・日時"),
        (MID_GRAY,"ImportBatch",     "作業計画の\nCSVインポート履歴"),
    ]

    tw, th = Inches(2.35), Inches(0.78)
    ml, mt = Inches(0.45), Inches(1.25)
    gap_x, gap_y = Inches(0.25), Inches(0.2)
    cols = 4

    for i, (color, name, desc) in enumerate(tables):
        col = i % cols
        row = i // cols
        l = ml + col * (tw + gap_x)
        t = mt + row * (th + gap_y)
        rect(sl, l, t, tw, th, fill=LIGHT_GRAY, lc=color, lw=0.75)
        lrect(sl, name, l, t, tw, Inches(0.32), fill=color, tc=WHITE, fs=10, bold=True)
        mlbox(sl, desc, l + Inches(0.1), t + Inches(0.35),
              tw - Inches(0.15), Inches(0.38), fs=9, color=DARK_GRAY)

    # データフロー
    rect(sl, Inches(0.45), Inches(5.5), Inches(12.35), Inches(0.85),
         fill=LIGHT_BLUE, lc=NAVY, lw=0.75)
    mlbox(sl, "データの流れ：　"
              "ImportBatch（計画インポート）"
              "  →  WorkPlan（作業計画）"
              "  →  WorkLog（作業実績）"
              "  →  InspectionReport + Defect（点検・不具合記録）"
              "  →  RepairLog（補修）",
          Inches(0.65), Inches(5.6), Inches(12.0), Inches(0.65),
          fs=10, color=NAVY, valign=MSO_ANCHOR.MIDDLE)

    mlbox(sl, "※ すべてのレコードはUser（作業者）と紐付けられ、誰がいつ何をしたかを追跡可能",
          Inches(0.45), Inches(6.5), Inches(12.3), Inches(0.35),
          fs=9, color=RGBColor(0x66, 0x66, 0x66))


# ============================================================
# メイン
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    s1(prs);  s2(prs);  s3(prs);  s4(prs)
    s5(prs);  s6(prs);  s7(prs);  s8(prs)
    s9(prs);  s10(prs); s11(prs)

    prs.save(OUT)
    print(f"保存完了: {OUT}")


if __name__ == "__main__":
    main()
