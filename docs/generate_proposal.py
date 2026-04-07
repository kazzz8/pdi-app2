#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
現場作業管理システム（仮称）上長向け提案書
PowerPoint生成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ============================================================
# 定数
# ============================================================
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "proposal_bucho.pptx")

NAVY       = RGBColor(0x1A, 0x3A, 0x6B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
LIGHT_BLUE = RGBColor(0xEA, 0xF0, 0xFB)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF8)
MID_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
YELLOW_BG  = RGBColor(0xFF, 0xFB, 0xE6)
ORANGE     = RGBColor(0xE8, 0xB0, 0x00)
FONT       = "游ゴシック"
W = Inches(13.33)
H = Inches(7.5)


# ============================================================
# ヘルパー
# ============================================================

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill=None, lc=None, lw=None):
    """矩形を追加。fill=Noneで透明背景。"""
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = lc
        if lw:
            s.line.width = Pt(lw)
    else:
        s.line.color.rgb = fill or WHITE
    return s


def set_para(tf, text, fs=11, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, first=True):
    """テキストフレームに段落を追加（first=Trueなら先頭段落を使用）。"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color
    return p


def multiline(slide, text, l, t, w, h,
              fs=11, bold=False, color=DARK_GRAY,
              align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    """改行付きテキストボックス（\\nで段落分割）。"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        set_para(tf, line, fs=fs, bold=bold, color=color,
                 align=align, first=(i == 0))
    return box


def label_rect(slide, text, l, t, w, h, fill, tc=WHITE,
               fs=12, bold=True, align=PP_ALIGN.CENTER,
               valign=MSO_ANCHOR.MIDDLE, lc=None):
    """背景色付き矩形＋テキスト。"""
    s = rect(slide, l, t, w, h, fill=fill, lc=lc or fill)
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    set_para(tf, text, fs=fs, bold=bold, color=tc, align=align)
    return s


def title_bar(slide, text):
    """ネイビーのタイトルバー（スライド上部）。"""
    bar = label_rect(slide, text,
                     Inches(0.4), Inches(0.3),
                     Inches(12.5), Inches(0.65),
                     fill=NAVY, tc=WHITE, fs=17, bold=True)
    return bar


# ============================================================
# スライド1: タイトル
# ============================================================
def s1_title(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=NAVY)
    # アクセントライン
    rect(sl, Inches(0.55), Inches(2.6), Inches(0.1), Inches(2.3), fill=WHITE)
    # タイトル
    multiline(sl, "現場作業管理システム（仮称）\n導入提案",
              Inches(0.85), Inches(2.3), Inches(11.5), Inches(1.8),
              fs=27, bold=True, color=WHITE)
    # サブタイトル
    multiline(sl, "現場の二重入力・紙作業を解消し、作業記録を一元管理する",
              Inches(0.85), Inches(4.25), Inches(11.5), Inches(0.6),
              fs=13, color=RGBColor(0xBB, 0xCC, 0xEE))
    # 日付
    multiline(sl, "2026年4月　情報システム部",
              Inches(0.85), Inches(5.1), Inches(6), Inches(0.5),
              fs=11, color=RGBColor(0x99, 0xAA, 0xCC))


# ============================================================
# スライド2: 現状の課題
# ============================================================
def s2_issues(sl_parent):
    sl = blank_slide(sl_parent)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "１．現状の課題")

    cards = [
        ("課題①\n三重の入力負担",
         "作業完了・不具合の記録を\n「紙」「自社システム」「マツダシステム」\nの3箇所に個別入力している。\n同じ情報を何度も入力する無駄が生じている。"),
        ("課題②\n共有PCへの移動",
         "作業時間の報告には共有PCを\n使用するため、作業場所から離れて\nPC前に並ぶ時間が毎日発生している。"),
        ("課題③\n自社システムの機能未達",
         "自社ベンダーシステムはマツダシステム\nとの自動連携を目的に導入されたが、\n連携は現在も未実現のまま。\n投資効果が十分に出ていない。"),
    ]

    cw, ch = Inches(3.9), Inches(4.6)
    lefts = [Inches(0.45), Inches(4.7), Inches(8.95)]
    top = Inches(1.35)

    for i, (title, body) in enumerate(cards):
        l = lefts[i]
        rect(sl, l, top, cw, ch,
             fill=RGBColor(0xFD, 0xFD, 0xFD), lc=MID_GRAY, lw=0.5)
        rect(sl, l, top, cw, Inches(0.06), fill=RED)
        label_rect(sl, title, l, top + Inches(0.06), cw, Inches(0.9),
                   fill=RGBColor(0xFF, 0xF0, 0xF0), tc=RED,
                   fs=12, bold=True, valign=MSO_ANCHOR.MIDDLE)
        multiline(sl, body,
                  l + Inches(0.15), top + Inches(1.08),
                  cw - Inches(0.3), ch - Inches(1.18),
                  fs=10.5, color=DARK_GRAY)

    multiline(sl, "※ 上記課題は作業者 200〜300名 全員に影響している",
              Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
              fs=9, color=RGBColor(0x88, 0x88, 0x88))


# ============================================================
# スライド3: 提案内容
# ============================================================
def s3_solution(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "２．提案内容")

    label_rect(sl, "現場作業管理システム（仮称）の導入",
               Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.85),
               fill=NAVY, tc=WHITE, fs=16, bold=True)

    features = [
        ("① スマートフォン・タブレットから\n　 その場で入力",
         "共有PCへの移動が不要。\n作業場所で即時記録が可能になる。"),
        ("② 作業記録を1システムに一元化",
         "作業時間・完了報告・不具合情報を\n1回の入力で完結する。"),
        ("③ 不具合写真のその場撮影・記録",
         "紙への記載が不要。\n写真とともに正確に記録できる。"),
        ("④ 情報システム部による内製開発",
         "追加のベンダーコストは発生しない。\n現場ニーズに合わせた柔軟な改修が可能。"),
    ]

    fw, fh = Inches(5.9), Inches(2.1)
    positions = [
        (Inches(0.5), Inches(2.25)),
        (Inches(6.9), Inches(2.25)),
        (Inches(0.5), Inches(4.55)),
        (Inches(6.9), Inches(4.55)),
    ]

    for i, (title, body) in enumerate(features):
        l, t = positions[i]
        rect(sl, l, t, fw, fh, fill=LIGHT_BLUE, lc=NAVY, lw=0.75)
        rect(sl, l, t, Inches(0.45), fh, fill=NAVY)
        multiline(sl, str(i + 1), l + Inches(0.0), t,
                  Inches(0.45), fh, fs=20, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        multiline(sl, title,
                  l + Inches(0.55), t + Inches(0.12),
                  fw - Inches(0.65), Inches(0.75),
                  fs=11, bold=True, color=NAVY)
        multiline(sl, body,
                  l + Inches(0.55), t + Inches(0.95),
                  fw - Inches(0.65), Inches(1.0),
                  fs=10, color=DARK_GRAY)


# ============================================================
# スライド4: 期待される効果
# ============================================================
def s4_effects(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "３．期待される効果")

    rows = [
        ("工数削減",
         "1人あたり1日 約5〜10分 の入力時間を削減",
         "200名換算で年間 約1,700〜3,300時間 の削減試算",
         RED, RGBColor(0xFD, 0xEC, 0xEA)),
        ("品質向上",
         "紙とシステムの転記ミス・記入漏れを防止",
         "不具合記録の抜け漏れが減少し、品質管理が向上",
         RED, RGBColor(0xFD, 0xEC, 0xEA)),
        ("コスト削減",
         "自社ベンダーシステムの保守費用を削減",
         "段階的移行後に廃止することで維持コストをゼロに",
         RED, RGBColor(0xFD, 0xEC, 0xEA)),
        ("即時性",
         "管理者がリアルタイムで作業進捗を把握可能",
         "集計・確認のための問い合わせ工数も削減",
         RGBColor(0xB7, 0x77, 0x0D), RGBColor(0xFE, 0xF9, 0xE7)),
    ]

    col_w = [Inches(1.9), Inches(5.6), Inches(4.8)]
    rh = Inches(1.05)
    hh = Inches(0.55)
    tl = Inches(0.5)
    tt = Inches(1.25)

    # ヘッダー
    for j, (hdr, cw) in enumerate(zip(["区分", "効果", "補足"], col_w)):
        x = tl + sum(col_w[:j])
        label_rect(sl, hdr, x, tt, cw, hh,
                   fill=NAVY, tc=WHITE, fs=11, bold=True)

    # データ行
    for i, (cat, eff, note, bc, bbg) in enumerate(rows):
        y = tt + hh + i * rh
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        x = tl

        # 区分セル
        rect(sl, x, y, col_w[0], rh, fill=bg, lc=MID_GRAY, lw=0.25)
        label_rect(sl, cat,
                   x + Inches(0.15), y + Inches(0.3),
                   col_w[0] - Inches(0.3), Inches(0.45),
                   fill=bbg, tc=bc, fs=10, bold=True, lc=bc)
        x += col_w[0]

        # 効果セル
        rect(sl, x, y, col_w[1], rh, fill=bg, lc=MID_GRAY, lw=0.25)
        multiline(sl, eff,
                  x + Inches(0.12), y + Inches(0.12),
                  col_w[1] - Inches(0.2), rh - Inches(0.2),
                  fs=10.5, bold=(i == 0), color=DARK_GRAY,
                  valign=MSO_ANCHOR.MIDDLE)
        x += col_w[1]

        # 補足セル
        rect(sl, x, y, col_w[2], rh, fill=bg, lc=MID_GRAY, lw=0.25)
        multiline(sl, note,
                  x + Inches(0.12), y + Inches(0.12),
                  col_w[2] - Inches(0.2), rh - Inches(0.2),
                  fs=10, color=DARK_GRAY, valign=MSO_ANCHOR.MIDDLE)

    multiline(sl, "※ 試算は1人5〜10分/日 × 200名 × 250日で算出。実際の削減効果は試験運用にて計測予定。",
              Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.4),
              fs=9, color=RGBColor(0x88, 0x88, 0x88))


# ============================================================
# スライド5: 既存システムとの関係
# ============================================================
def s5_systems(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "４．既存システムとの関係")

    bh = Inches(4.6)
    bt = Inches(1.3)

    # --- 現状 ---
    rect(sl, Inches(0.4), bt, Inches(5.4), bh,
         fill=LIGHT_GRAY, lc=MID_GRAY, lw=0.75)
    label_rect(sl, "現状（3システム並存）",
               Inches(0.4), bt, Inches(5.4), Inches(0.55),
               fill=RGBColor(0x77, 0x77, 0x77), tc=WHITE, fs=12)

    curr = [
        ("マツダシステム",          RGBColor(0x22, 0x55, 0xAA),
         "完成車輸送管理・作業完了報告"),
        ("自社ベンダーシステム",    RGBColor(0x77, 0x77, 0x77),
         "作業時間管理・作業完了報告\nマツダ連携：未実現のまま"),
        ("紙（点検票・不具合票）",  RGBColor(0x99, 0x66, 0x00),
         "各種記録・報告"),
    ]
    for i, (name, color, desc) in enumerate(curr):
        it = bt + Inches(0.65) + i * Inches(1.25)
        rect(sl, Inches(0.6), it, Inches(5.0), Inches(1.1),
             fill=WHITE, lc=color, lw=1.0)
        multiline(sl, name,
                  Inches(0.75), it + Inches(0.1),
                  Inches(4.8), Inches(0.4),
                  fs=11, bold=True, color=color)
        multiline(sl, desc,
                  Inches(0.75), it + Inches(0.55),
                  Inches(4.8), Inches(0.5),
                  fs=9.5, color=DARK_GRAY)

    # 矢印
    multiline(sl, "→",
              Inches(6.0), Inches(3.5), Inches(1.2), Inches(0.9),
              fs=38, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # --- 移行後 ---
    rect(sl, Inches(7.4), bt, Inches(5.5), bh,
         fill=RGBColor(0xEA, 0xF0, 0xFB), lc=NAVY, lw=1.0)
    label_rect(sl, "移行後（役割を整理）",
               Inches(7.4), bt, Inches(5.5), Inches(0.55),
               fill=NAVY, tc=WHITE, fs=12)

    future = [
        ("マツダシステム",                 GREEN,
         "継続利用（変更なし）"),
        ("現場作業管理システム（仮称）",   NAVY,
         "作業時間・完了・不具合を一元管理\n自社ベンダーシステムを置き換え"),
        ("紙：廃止",                       RED,
         ""),
    ]
    for i, (name, color, desc) in enumerate(future):
        it = bt + Inches(0.65) + i * Inches(1.25)
        rect(sl, Inches(7.6), it, Inches(5.1), Inches(1.1),
             fill=WHITE, lc=color, lw=1.0)
        multiline(sl, name,
                  Inches(7.75), it + Inches(0.1),
                  Inches(4.8), Inches(0.4),
                  fs=11, bold=True, color=color)
        if desc:
            multiline(sl, desc,
                      Inches(7.75), it + Inches(0.55),
                      Inches(4.8), Inches(0.5),
                      fs=9.5, color=DARK_GRAY)

    multiline(sl, "※ マツダシステムへの報告は従来通り作業者が行う。将来的な自動連携は中長期の検討事項とする。",
              Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
              fs=9, color=RGBColor(0x88, 0x88, 0x88))


# ============================================================
# スライド6: 導入スケジュール
# ============================================================
def s6_schedule(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "５．導入スケジュール（案）")

    phases = [
        ("Phase 1", "〜2026年夏頃",
         "・開発完了\n・社内テスト\n・インフラ準備"),
        ("Phase 2", "2026年秋頃",
         "・1〜2チームで\n　試験運用開始\n・現場フィードバック収集"),
        ("Phase 3", "2026年度内",
         "・全体展開\n・紙運用廃止\n・効果測定"),
        ("Phase 4", "移行完了後",
         "・自社ベンダー\n　システム廃止\n・保守費用削減"),
    ]

    pw, ph = Inches(2.85), Inches(4.0)
    pt = Inches(1.5)
    lefts = [Inches(0.45), Inches(3.6), Inches(6.75), Inches(9.9)]

    for i, (phase, period, content) in enumerate(phases):
        l = lefts[i]
        rect(sl, l, pt, pw, ph, fill=WHITE, lc=MID_GRAY, lw=0.5)
        label_rect(sl, phase, l, pt, pw, Inches(0.55),
                   fill=NAVY, tc=WHITE, fs=14, bold=True)
        multiline(sl, period,
                  l, pt + Inches(0.6), pw, Inches(0.45),
                  fs=10, color=RGBColor(0x66, 0x66, 0x66),
                  align=PP_ALIGN.CENTER)
        rect(sl, l + Inches(0.2), pt + Inches(1.1),
             pw - Inches(0.4), Inches(0.02), fill=MID_GRAY)
        multiline(sl, content,
                  l + Inches(0.2), pt + Inches(1.2),
                  pw - Inches(0.3), ph - Inches(1.3),
                  fs=11, color=DARK_GRAY)

        if i < 3:
            multiline(sl, "▶",
                      l + pw + Inches(0.02), pt + ph / 2 - Inches(0.3),
                      Inches(0.5), Inches(0.55),
                      fs=14, bold=True, color=NAVY,
                      align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    multiline(sl, "※ スケジュールは目安。セキュリティ審査・稟議の進行状況により変動あり。",
              Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
              fs=9, color=RGBColor(0x88, 0x88, 0x88))


# ============================================================
# スライド7: ご確認・ご支援のお願い
# ============================================================
def s7_action(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=WHITE)
    title_bar(sl, "６．ご確認・ご支援のお願い")

    multiline(sl, "本提案の推進にあたり、以下についてご意見・ご支援をいただけますと幸いです。",
              Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.45),
              fs=11.5, color=DARK_GRAY)

    actions = [
        ("① 提案の方向性へのご賛同",
         "現状の課題認識と、本システムによる解決方向性についてご意見をいただけますか。"),
        ("② 試験運用チームのご紹介",
         "ご協力いただける現場チームへのご紹介・ご調整をお願いできますでしょうか。"),
        ("③ 社内調整へのご支援",
         "情報システム部・セキュリティ担当との調整において、後ろ盾としてのご支援をお願いできますでしょうか。"),
    ]

    for i, (title, body) in enumerate(actions):
        it = Inches(1.9) + i * Inches(1.5)
        rect(sl, Inches(0.5), it, Inches(12.3), Inches(1.3),
             fill=YELLOW_BG, lc=ORANGE, lw=0.75)
        multiline(sl, title,
                  Inches(0.7), it + Inches(0.1),
                  Inches(12.0), Inches(0.5),
                  fs=12, bold=True, color=RGBColor(0x7D, 0x5A, 0x00))
        multiline(sl, body,
                  Inches(0.7), it + Inches(0.65),
                  Inches(12.0), Inches(0.55),
                  fs=10.5, color=DARK_GRAY)

    multiline(sl, "情報システム部　2026年4月",
              Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.45),
              fs=10, color=RGBColor(0x88, 0x88, 0x88),
              align=PP_ALIGN.RIGHT)


# ============================================================
# メイン
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    s1_title(prs)
    s2_issues(prs)
    s3_solution(prs)
    s4_effects(prs)
    s5_systems(prs)
    s6_schedule(prs)
    s7_action(prs)

    prs.save(OUT)
    print(f"保存完了: {OUT}")


if __name__ == "__main__":
    main()
